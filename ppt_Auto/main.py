from PySide6.QtWidgets import QApplication, QMainWindow, QFileDialog, QMessageBox
from PySide6.QtGui import QIcon
from mainWindow import Ui_MainWindow
from pptx import Presentation
from datetime import datetime
import PIL
from PIL import ImageGrab, Image
import pyautogui
import icon

#아직 구현하지 못한 것.
# 1. clean code
# 2. 다중 슬라이드 생성 및 편집.
# 3. 이미지 파일 사용 후 제거?
# 4. 외관? -> self.setWindowIcon() 해서 로고 이미지 추가.

class mainWindow(QMainWindow, Ui_MainWindow):
    def __init__(self):
        super(mainWindow, self).__init__()
        self.setupUi(self)
        self.init_ui()
        self.show_ver.triggered.connect(self.show_inform)
        self.close_window.triggered.connect(self.close_all)

    # show logo image
    def init_ui(self):
        self.setWindowTitle('AutoPPT')
        self.setWindowIcon(QIcon(':/logo.png'))
        # setGeomety (X, Y, Width, Height)
        self.setGeometry(400, 100, 500, 500)
        self.show()

    def close_all(self):
        QMessageBox.about(self, "End", "BYE!")
        QApplication.closeAllWindows()

    def show_inform(self):
        QMessageBox.about(self, "Version", "Demo 0.0.3r")

    # make ppt - This is a sample code.
    # Input your ppt code here.
    def make_ppt(self):
        try:
            # your code Start --------------------------------------
            title = self.titleEdit.text()
            inform = self.informEdit.toPlainText()
            now = datetime.now()
            current_time = now.strftime("%Y-%m-%d %H:%M:%S\n")
            current_text = current_time + inform
            img_path = self.imgLabel.text()

            prs = Presentation('base.pptx')
            main_slide = prs.slides[0]

            shape_list = main_slide.shapes
            shape_index = {}
            for i, shape in enumerate(shape_list):
                shape_index[shape.name] = i

            # ppt Title
            title_select = shape_list[shape_index['title_box']].text_frame
            title_select.text = title

            # ppt Image
            img_select = shape_list[shape_index["img_box"]]
            img_select.insert_picture(img_path)

            # ppt Main-text
            text_select = shape_list[shape_index['text_box']].text_frame
            text_select.text = current_text
            
            # Setting output
            output = now.strftime('%Y-%m-%d %H_%M') + '.pptx'
            prs.save(output)
            QMessageBox.about(self, "About", "Work Done!")
            QApplication.closeAllWindows()
            # your code end----------------------------------------------------
        # Setting Exception
        except FileNotFoundError:
            QMessageBox.warning(self, "Error!", "There is no image File.\n Please select image file or capture screen.")
        except PIL.UnidentifiedImageError:
            QMessageBox.warning(self, "Error!", "You selected wrong type of file.\n File type is must be image.")

    def capture_screen(self):
        pyautogui.hotkey('win', 'shift', 's')

    def upload_func(self):
        try:
            img_file = QFileDialog.getOpenFileName(self, filter="PNG (*.png);; JPEG (*.jpg *.jpeg);; All files (*.*)")
            self.imgLabel.setText(img_file[0])
        except IndexError:
            pass

    def get_img(self):
        now = datetime.now()
        current_time = now.strftime("%Y-%m-%d %H_%M_%S")
        img_name = current_time + '.png'
        # get image from clipboard
        img = ImageGrab.grabclipboard()
        img.save(img_name, 'PNG')
        self.imgLabel.setText(img_name)
        # resize image
        im = Image.open(img_name)
        im = im.resize((644, 400))
        im = im.convert('RGB')
        im.save(img_name, 'PNG')


app = QApplication()
win = mainWindow()
win.show()
app.exec()
